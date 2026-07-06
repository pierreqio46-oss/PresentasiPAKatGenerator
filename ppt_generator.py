from pptx import Presentation
from pptx.util import Inches

class PPTGenerator:

    def __init__(
            self,
            image_generator):

        self.image_generator=(
            image_generator
        )


    def create(
            self,
            content,
            quiz):

        prs=Presentation()

        cover=prs.slides.add_slide(
            prs.slide_layouts[0]
        )

        cover.shapes.title.text=(
            "Presentasi Modul Ajar"
        )

        slides=content.split(
            "Slide"
        )

        for item in slides:

            if len(item)<10:
                continue

            slide=prs.slides.add_slide(
                prs.slide_layouts[5]
            )

            lines=item.split(
                "\n"
            )

            slide.shapes.title.text=(
                lines[0]
            )

            box=slide.shapes.add_textbox(
                Inches(0.5),
                Inches(1.3),
                Inches(6),
                Inches(4)
            )

            frame=box.text_frame

            for line in lines[1:]:

                if line.strip():

                    p=frame.add_paragraph()

                    p.text="• "+line


            try:

                img=(
                    self.image_generator
                    .create(lines[0])
                )

                slide.shapes.add_picture(
                    img,
                    Inches(7),
                    Inches(1.5),
                    width=Inches(3)
                )

            except:
                pass


        quiz_slide=prs.slides.add_slide(
            prs.slide_layouts[1]
        )

        quiz_slide.shapes.title.text=(
            "Kuis"
        )

        quiz_slide.placeholders[
            1
        ].text=quiz

        output=(
            "output/presentasi.pptx"
        )

        prs.save(output)

        return output